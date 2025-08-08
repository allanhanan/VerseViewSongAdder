import sys
import os
import sqlite3
import shutil
from datetime import datetime
from pptx import Presentation
import win32com.client
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QLabel, QPushButton,
    QFileDialog, QListWidget, QProgressBar, QMessageBox, QHBoxLayout,
    QGroupBox, QDialog, QFormLayout, QLineEdit,
    QListWidgetItem, QAbstractItemView, QTextEdit, QStyle, QStyleOptionButton,
    QStyledItemDelegate
)
from PyQt5.QtCore import Qt, QSize, QRect
from PyQt5.QtGui import QPainter, QMouseEvent, QIcon
import qdarkstyle
import getpass
import glob

# Custom Dialog for Settings
class SettingsDialog(QDialog):
    """A dialog to customize default font and category for song imports."""
    def __init__(self, parent=None, font="Calibri", category="autoadd"):
        super().__init__(parent)
        self.setWindowTitle("Customize Settings")
        self.setStyleSheet(qdarkstyle.load_stylesheet_pyqt5())
        self.setFixedSize(400, 150)
        
        layout = QFormLayout()
        
        # User can set a default font.
        self.font_input = QLineEdit(self)
        self.font_input.setText(font)
        layout.addRow("Default Font:", self.font_input)
        
        # User can set a default category.
        self.category_input = QLineEdit(self)
        self.category_input.setText(category)
        layout.addRow("Default Category:", self.category_input)
        
        button_layout = QHBoxLayout()
        ok_button = QPushButton("OK")
        cancel_button = QPushButton("Cancel")
        
        ok_button.clicked.connect(self.accept)
        cancel_button.clicked.connect(self.reject)
        
        button_layout.addWidget(ok_button)
        button_layout.addWidget(cancel_button)
        
        layout.addRow(button_layout)
        self.setLayout(layout)

    def get_settings(self):
        """Returns the font and category entered by the user."""
        return self.font_input.text(), self.category_input.text()

# Custom delegate to draw delete buttons
class ButtonDelegate(QStyledItemDelegate):
    """Paints a small 'x' button on each list item for individual deletion."""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.parent_widget = parent

    def paint(self, painter: QPainter, option, index):
        """Paints the button and list item text."""
        super().paint(painter, option, index)
        
        # Calculate the button's position.
        button_rect = QRect(
            option.rect.right() - 25,
            option.rect.top() + (option.rect.height() - 20) // 2,
            20, 20
        )
        
        button_option = QStyleOptionButton()
        button_option.rect = button_rect
        button_option.text = "x"
        button_option.state = QStyle.State_Enabled
        
        # Change button color based on selection state for better visibility.
        if option.state & QStyle.State_Selected:
            button_option.palette.setColor(button_option.palette.ButtonText, Qt.white)
        else:
            button_option.palette.setColor(button_option.palette.ButtonText, Qt.red)
        
        QApplication.style().drawControl(QStyle.CE_PushButton, button_option, painter, None)

    def editorEvent(self, event, model, option, index):
        """Handles clicks on the 'x' button to delete the corresponding file."""
        if event.type() == QMouseEvent.MouseButtonPress:
            button_rect = QRect(
                option.rect.right() - 25,
                option.rect.top() + (option.rect.height() - 20) // 2,
                20, 20
            )
            
            if button_rect.contains(event.pos()):
                # Delete the file when the 'x' button is clicked.
                self.parent_widget.delete_single_file(index.row())
                return True
        return super().editorEvent(event, model, option, index)

# Custom ListWidget to handle drag-and-drop and key press events
class CustomListWidget(QListWidget):
    """Extends QListWidget for drag-and-drop and delete key functionality."""
    def __init__(self, main_window, parent=None):
        super().__init__(parent)
        self.main_window = main_window
        self.setSelectionMode(QAbstractItemView.ExtendedSelection)
        self.setAcceptDrops(True)
        self.setDragDropMode(QAbstractItemView.InternalMove)
        self.setAlternatingRowColors(True)
        self.setItemDelegate(ButtonDelegate(self.main_window))

    def keyPressEvent(self, event):
        """Triggers deletion of selected items when the delete key is pressed."""
        if event.key() == Qt.Key_Delete:
            self.main_window.delete_selected()
        super().keyPressEvent(event)

    def dragEnterEvent(self, event):
        """Allows drag-and-drop of files from outside the application."""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
    
    def dropEvent(self, event):
        """Adds dropped files/folders to the file list."""
        files_to_add = []
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if os.path.isdir(path):
                # Recursively add files from dropped folders.
                for f in os.listdir(path):
                    full_path = os.path.join(path, f)
                    if full_path.lower().endswith((".pptx", ".ppt")):
                        files_to_add.append(full_path)
            elif path.lower().endswith((".pptx", ".ppt")):
                files_to_add.append(path)
        self.main_window.add_files_to_list(files_to_add)
        event.accept()

# Main Application Window
class SongDBInjector(QWidget):
    """The main application window for injecting songs into a VerseVIEW database."""
    def __init__(self):
        super().__init__()
        self.setWindowTitle("VerseView Song Adder")
        self.resize(1200, 800)
        self.setStyleSheet(qdarkstyle.load_stylesheet_pyqt5())
        self.setWindowIcon(QIcon("app_icon.ico"))
        self.db_path = None
        self.files = []
        self.layout_widgets()
        self.auto_find_db()
        self.setup_connections()

    def layout_widgets(self):
        """Sets up the UI elements and their layout."""
        main_layout = QVBoxLayout()
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        title_label = QLabel("VerseView Song Adder\n Add Songs from PPT's Automatically")
        title_label.setAlignment(Qt.AlignCenter)
        title_label.setStyleSheet("font-size: 32px; font-weight: bold; color: #E0E0E0; margin-bottom: 10px;")
        main_layout.addWidget(title_label)
        
        # Main split layout.
        content_layout = QHBoxLayout()
        content_layout.setSpacing(20)

        # Left side: DB and File management.
        left_layout = QVBoxLayout()
        left_layout.setSpacing(15)
        
        # Database Selection section.
        db_group = QGroupBox("Database Selection")
        db_group.setStyleSheet("font-size: 18px; font-weight: bold;")
        db_layout = QVBoxLayout()
        self.db_label = QLabel("Looking for songs.db...")
        self.db_label.setWordWrap(True)
        self.db_label.setStyleSheet("font-size: 16px; font-weight: normal; margin-bottom: 5px;")
        db_layout.addWidget(self.db_label)
        
        db_button = QPushButton("Choose .db File Manually")
        db_button.setFixedHeight(40)
        db_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        db_button.setToolTip("Click to manually select the songs.db file if it was not found automatically.")
        db_layout.addWidget(db_button)
        self.db_button = db_button
        db_group.setLayout(db_layout)
        left_layout.addWidget(db_group)
        
        # File Selection section.
        file_group = QGroupBox("Song Files (.ppt / .pptx)")
        file_group.setStyleSheet("font-size: 18px; font-weight: bold;")
        file_layout = QVBoxLayout()

        file_buttons_layout = QHBoxLayout()
        file_buttons_layout.setSpacing(10)
        self.folder_button = QPushButton("Scan Folder")
        self.folder_button.setFixedHeight(40)
        self.folder_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        self.folder_button.setToolTip("Select a folder to automatically add all PowerPoint files from it.")
        file_buttons_layout.addWidget(self.folder_button)

        self.add_file_button = QPushButton("Add File(s)")
        self.add_file_button.setFixedHeight(40)
        self.add_file_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        self.add_file_button.setToolTip("Select one or more PowerPoint files to add to the list.")
        file_buttons_layout.addWidget(self.add_file_button)
        
        self.clear_list_button = QPushButton("Clear List")
        self.clear_list_button.setFixedHeight(40)
        self.clear_list_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        self.clear_list_button.setToolTip("Remove all files from the list.")
        file_buttons_layout.addWidget(self.clear_list_button)
        
        self.delete_selected_button = QPushButton("Delete Selected")
        self.delete_selected_button.setFixedHeight(40)
        self.delete_selected_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        self.delete_selected_button.setToolTip("Remove selected files from the list. Use Shift or Ctrl to select multiple files.")
        file_buttons_layout.addWidget(self.delete_selected_button)

        file_layout.addLayout(file_buttons_layout)

        self.file_list = CustomListWidget(self)
        self.file_list.setStyleSheet("font-size: 14px; padding: 5px;")
        self.file_list.setToolTip("Drag and drop PowerPoint files here. You can also use the Delete key to remove selected files.")
        file_layout.addWidget(self.file_list)
        file_group.setLayout(file_layout)
        left_layout.addWidget(file_group)
        
        content_layout.addLayout(left_layout, 2)

        # Right side: Preview and Controls.
        right_layout = QVBoxLayout()
        right_layout.setSpacing(15)

        # Preview pane section.
        preview_group = QGroupBox("Song Preview")
        preview_group.setStyleSheet("font-size: 18px; font-weight: bold;")
        preview_layout = QVBoxLayout()
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        self.preview_text.setStyleSheet("font-size: 14px; padding: 10px; border: 1px solid #555;")
        preview_layout.addWidget(self.preview_text)
        preview_group.setLayout(preview_layout)
        right_layout.addWidget(preview_group, 1)

        # Control Panel section.
        controls_group = QGroupBox("Settings & Actions")
        controls_group.setStyleSheet("font-size: 18px; font-weight: bold;")
        controls_layout = QVBoxLayout()
        
        self.custom_font_label = QLabel("Font: Calibri")
        self.custom_cat_label = QLabel("Category: autoadd")
        self.custom_font_label.setStyleSheet("font-size: 14px;")
        self.custom_cat_label.setStyleSheet("font-size: 14px;")
        controls_layout.addWidget(self.custom_font_label)
        controls_layout.addWidget(self.custom_cat_label)

        customize_button = QPushButton("Customize Settings")
        customize_button.setFixedHeight(40)
        customize_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        customize_button.setToolTip("Change the default font and category for new songs.")
        self.customize_button = customize_button
        controls_layout.addWidget(customize_button)
        
        backup_button = QPushButton("Backup Database")
        backup_button.setFixedHeight(40)
        backup_button.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px; border-radius: 5px;")
        backup_button.setToolTip("Create a backup copy of the current songs.db before making changes.")
        self.backup_button = backup_button
        controls_layout.addWidget(backup_button)

        self.progress = QProgressBar()
        self.progress.setVisible(False)
        controls_layout.addWidget(self.progress)
        
        self.inject_button = QPushButton("Add Songs to Database")
        self.inject_button.setFixedHeight(50)
        self.inject_button.setStyleSheet("font-size: 18px; font-weight: bold; padding: 5px; border-radius: 5px; background-color: #2e8b57;")
        self.inject_button.setToolTip("Click to add all files from the list to the selected database.")
        controls_layout.addWidget(self.inject_button)

        controls_group.setLayout(controls_layout)
        controls_group.setFixedWidth(350)
        right_layout.addWidget(controls_group)
        
        content_layout.addLayout(right_layout, 1)
        
        main_layout.addLayout(content_layout)
        self.setLayout(main_layout)

        self.default_font = "Calibri"
        self.default_category = "autoadd"

    def setup_connections(self):
        """Connects UI elements to their corresponding functions."""
        self.db_button.clicked.connect(self.choose_db)
        self.folder_button.clicked.connect(self.scan_folder)
        self.add_file_button.clicked.connect(self.add_file)
        self.clear_list_button.clicked.connect(self.clear_list)
        self.delete_selected_button.clicked.connect(self.delete_selected)
        self.file_list.itemSelectionChanged.connect(self.preview_selected_file)
        self.customize_button.clicked.connect(self.customize_settings)
        self.backup_button.clicked.connect(self.backup_db)
        self.inject_button.clicked.connect(self.inject_all)

    def auto_find_db(self):
        """Attempts to automatically locate the VerseVIEW songs.db file."""
        user = getpass.getuser()
        base_path = f"C:/Users/{user}/AppData/Roaming/"
        search_pattern = os.path.join(base_path, "VerseVIEW*", "vvdata", "songs", "songs.db")
        found = glob.glob(search_pattern)
        if found:
            self.db_path = found[0]
            self.db_label.setText(f"Auto-selected DB: {self.db_path}")
            self.db_label.setToolTip(self.db_path)
        else:
            self.db_label.setText("songs.db not found automatically")
            self.db_label.setToolTip("Please click 'Choose .db File Manually' to select the database file.")

    def choose_db(self):
        """Opens a file dialog for manual database selection."""
        path, _ = QFileDialog.getOpenFileName(self, "Select Database", "", "*.db")
        if path:
            self.db_path = path
            self.db_label.setText(f"Selected DB: {path}")

    def scan_folder(self):
        """Scans a selected folder for PowerPoint files and adds them."""
        folder = QFileDialog.getExistingDirectory(self, "Select Folder")
        if folder:
            new_files = [os.path.join(folder, f) for f in os.listdir(folder) if f.lower().endswith((".pptx", ".ppt"))]
            self.add_files_to_list(new_files)

    def add_file(self):
        """Opens a file dialog for adding one or more PowerPoint files."""
        new_files, _ = QFileDialog.getOpenFileNames(self, "Add File(s)", "", "PowerPoint Files (*.pptx *.ppt)")
        if new_files:
            self.add_files_to_list(new_files)

    def add_files_to_list(self, new_files):
        """Adds a list of new files to the application's internal list and UI."""
        for file in new_files:
            if file not in self.files:
                self.files.append(file)
                self.file_list.addItem(os.path.basename(file))
        self.update_file_list()

    def clear_list(self):
        """Clears all files from the list with a confirmation."""
        if not self.files:
            return
        
        reply = QMessageBox.question(self, "Clear List",
                                     "Are you sure you want to clear all files from the list?",
                                     QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
        
        if reply == QMessageBox.Yes:
            self.files = []
            self.update_file_list()

    def delete_single_file(self, row):
        """Deletes a single file from the list without confirmation.
        
        This is triggered by the 'x' button on a list item.
        """
        if 0 <= row < len(self.files):
            del self.files[row]
            self.update_file_list()

    def delete_selected(self):
        """Deletes selected files, with confirmation for multiple files.
        
        This method handles both the 'Delete Selected' button and the keyboard shortcut.
        It bypasses the confirmation prompt for a single selected file.
        """
        selected_items = self.file_list.selectedItems()
        if not selected_items:
            return
        
        if len(selected_items) == 1:
            # Delete without confirmation when only one file is selected.
            item_to_delete = selected_items[0].text()
            self.files = [f for f in self.files if os.path.basename(f) != item_to_delete]
            self.update_file_list()
        else:
            # Show a confirmation prompt for multiple files.
            reply = QMessageBox.question(self, "Delete Files",
                                         f"Are you sure you want to delete {len(selected_items)} selected files?",
                                         QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
            
            if reply == QMessageBox.Yes:
                items_to_delete = [item.text() for item in selected_items]
                self.files = [f for f in self.files if os.path.basename(f) not in items_to_delete]
                self.update_file_list()

    def update_file_list(self):
        """Refreshes the file list widget to reflect changes."""
        self.file_list.clear()
        for file in self.files:
            item = QListWidgetItem(os.path.basename(file))
            item.setSizeHint(QSize(20, 25))
            # Added a tooltip to show the full file path on hover.
            item.setToolTip(file)
            self.file_list.addItem(item)
        if self.files:
            self.file_list.setCurrentRow(0)
        else:
            self.preview_text.clear()

    def preview_selected_file(self):
        """Previews the lyrics of the currently selected file."""
        selected_items = self.file_list.selectedItems()
        if not selected_items:
            self.preview_text.clear()
            return
        
        file_name = selected_items[0].text()
        file_path = next((f for f in self.files if os.path.basename(f) == file_name), None)
        
        if file_path:
            lyrics = self.extract_lyrics(file_path)
            formatted_lyrics = lyrics.replace("<slide>", "\n\n---\n\n").replace("<BR>", "\n")
            self.preview_text.setText(formatted_lyrics)

    def customize_settings(self):
        """Opens the settings dialog for custom configuration."""
        dialog = SettingsDialog(self, self.default_font, self.default_category)
        if dialog.exec_() == QDialog.Accepted:
            font, category = dialog.get_settings()
            self.default_font = font
            self.default_category = category
            self.custom_font_label.setText(f"Font: {self.default_font}")
            self.custom_cat_label.setText(f"Category: {self.default_category}")

    def backup_db(self):
        """Creates a timestamped backup of the selected database."""
        if not self.db_path:
            QMessageBox.warning(self, "Missing DB", "Please select a database first.")
            return

        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        backup_path = f"{self.db_path}.backup_{timestamp}"
        try:
            shutil.copyfile(self.db_path, backup_path)
            QMessageBox.information(self, "Backup Complete", f"Database backed up to:\n{backup_path}")
        except Exception as e:
            QMessageBox.critical(self, "Backup Failed", f"An error occurred during backup: {e}")

    def extract_text_pptx(self, path):
        """Extracts text from a .pptx file using the python-pptx library."""
        prs = Presentation(path)
        all_text = []
        for slide in prs.slides:
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for line in shape.text.splitlines():
                        line = line.rstrip()
                        if line:
                            lines.append(line)
            slide_block = "<BR>".join(lines)
            all_text.append(slide_block)
        return "<slide>".join(all_text)

    def extract_text_ppt(self, path):
        """Extracts text from an older .ppt file using the win32com library."""
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        presentation = ppt_app.Presentations.Open(path, WithWindow=False)
        all_text = []
        for slide in presentation.Slides:
            lines = []
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text
                    for line in text.splitlines():
                        line = line.strip()
                        if line:
                            lines.append(line)
            slide_block = "<BR>".join(lines)
            all_text.append(slide_block)
        presentation.Close()
        ppt_app.Quit()
        return "<slide>".join(all_text)

    def extract_lyrics(self, file_path):
        """Determines the file type and calls the appropriate text extraction method."""
        try:
            if file_path.lower().endswith(".pptx"):
                return self.extract_text_pptx(file_path)
            elif file_path.lower().endswith(".ppt"):
                return self.extract_text_ppt(file_path)
        except Exception as e:
            print(f"Error extracting {file_path}: {e}")
            return ""

    def get_next_id(self, conn):
        """Fetches the next available song ID from the database."""
        cur = conn.cursor()
        cur.execute("SELECT MAX(id) FROM sm")
        result = cur.fetchone()
        return (result[0] or 0) + 1

    def song_exists(self, conn, name):
        """Checks if a song with the given name already exists in the database."""
        cur = conn.cursor()
        cur.execute("SELECT id FROM sm WHERE name = ?", (name,))
        return cur.fetchone() is not None

    def extract_name_from_filename(self, file_path):
        """Extracts the song name from the filename."""
        base = os.path.basename(file_path)
        name, _ = os.path.splitext(base)
        return name.strip()

    def inject_all(self):
        """Iterates through all files and injects their content into the database."""
        if not self.db_path or not self.files:
            QMessageBox.warning(self, "Missing Info", "Database or files are not selected.")
            return

        confirm = QMessageBox.question(
            self, "Confirm Injection",
            f"Inject {len(self.files)} file(s) into the database?",
            QMessageBox.Yes | QMessageBox.No
        )
        if confirm != QMessageBox.Yes:
            return

        conn = sqlite3.connect(self.db_path)
        self.progress.setVisible(True)
        self.progress.setMaximum(len(self.files))
        self.progress.setValue(0)
        
        added_names = []
        failed_files = []
        
        for i, file in enumerate(self.files, start=1):
            lyrics = self.extract_lyrics(file)
            name = self.extract_name_from_filename(file)

            if not lyrics:
                failed_files.append(f"{name} (Error extracting lyrics)")
                self.progress.setValue(i)
                continue

            if self.song_exists(conn, name):
                reply = QMessageBox.question(
                    self, "Duplicate Song",
                    f"A song named '{name}' already exists. Overwrite it?",
                    QMessageBox.Yes | QMessageBox.No
                )
                if reply == QMessageBox.Yes:
                    cur = conn.cursor()
                    cur.execute("UPDATE sm SET lyrics = ? WHERE name = ?", (lyrics, name))
                    conn.commit()
                    added_names.append(f"{name} (Overwritten)")
                else:
                    failed_files.append(f"{name} (Skipped, duplicate)")
            else:
                new_id = self.get_next_id(conn)
                data = (
                    new_id, name, self.default_category, self.default_font, None, None,
                    "", "", "", "", "", lyrics,
                    "", "", "", "", 0, 0, 0, "null"
                )
                try:
                    cur = conn.cursor()
                    cur.execute('''INSERT INTO sm (
                        id, name, cat, font, font2, timestamp,
                        yvideo, bkgndfname, key, copy, notes,
                        lyrics, lyrics2, title2, tags,
                        slideseq, rating, chordsavailable, usagecount, subcat
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''', data)
                    conn.commit()
                    added_names.append(name)
                except Exception as e:
                    failed_files.append(f"{name} (DB Error: {e})")

            self.progress.setValue(i)
            QApplication.processEvents()

        conn.close()
        self.progress.setVisible(False)

        summary_message = "Injection Complete!\n\n"
        if added_names:
            summary_message += "Added/Updated Songs:\n" + "\n".join(added_names) + "\n\n"
        if failed_files:
            summary_message += "Failed Files:\n" + "\n".join(failed_files)
        
        QMessageBox.information(self, "Injection Summary", summary_message)
        self.files = []
        self.update_file_list()

def main():
    """Initializes and runs the application."""
    app = QApplication(sys.argv)
    win = SongDBInjector()
    win.show()
    sys.exit(app.exec_())

if __name__ == "__main__":
    main()